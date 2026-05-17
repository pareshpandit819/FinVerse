"""
Generates FinVerse demo presentation as a .pptx file.
Run: python3 docs/build_pptx.py
Output: docs/FinVerse_Demo_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# ── colour palette ────────────────────────────────────────────────────────────
SKY_950  = RGBColor(0x08, 0x2f, 0x49)
SKY_900  = RGBColor(0x0c, 0x4a, 0x6e)
SKY_700  = RGBColor(0x03, 0x69, 0xa1)
SKY_500  = RGBColor(0x0e, 0xa5, 0xe9)
SKY_400  = RGBColor(0x38, 0xbd, 0xf8)
SKY_100  = RGBColor(0xe0, 0xf2, 0xfe)
SKY_50   = RGBColor(0xf0, 0xf9, 0xff)
WHITE    = RGBColor(0xff, 0xff, 0xff)
GRAY_700 = RGBColor(0x37, 0x41, 0x51)
GRAY_500 = RGBColor(0x6b, 0x72, 0x80)
GRAY_200 = RGBColor(0xe5, 0xe7, 0xeb)
ROSE_500 = RGBColor(0xf4, 0x3f, 0x5e)
ROSE_100 = RGBColor(0xff, 0xe4, 0xe6)
EMERALD  = RGBColor(0x05, 0x96, 0x69)
EMLD_100 = RGBColor(0xd1, 0xfa, 0xe5)
AMBER    = RGBColor(0xf5, 0x9e, 0x0b)
AMBER_100= RGBColor(0xfe, 0xf3, 0xc7)
VIOLET   = RGBColor(0x7c, 0x3a, 0xed)
VIOLT_100= RGBColor(0xed, 0xe9, 0xfe)

# ── slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # completely blank

# ── helpers ────────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank)

def rect(slide, x, y, w, h, fill_rgb=None, fill_alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    return shape

def textbox(slide, x, y, w, h, text, size=14, bold=False, color=GRAY_700,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def multiline(slide, x, y, w, h, lines, base_size=12, color=GRAY_700):
    """lines = list of (text, size, bold, color) tuples"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, clr) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = clr or color
        run.font.name = "Calibri"

def label_tag(slide, x, y, text):
    """Small uppercase label in sky-500."""
    textbox(slide, x, y, 6, 0.3, text, size=9, bold=True,
            color=SKY_500, align=PP_ALIGN.LEFT)

def section_card(slide, x, y, w, h, title, body, title_color=SKY_950, bg=SKY_50, border=SKY_100):
    r = rect(slide, x, y, w, h, fill_rgb=bg)
    # border via outline
    r.line.color.rgb = border
    r.line.width = Pt(0.75)
    # round corners not supported in pptx easily — use rectangle
    textbox(slide, x+0.15, y+0.12, w-0.3, 0.28, title, size=10, bold=True, color=title_color)
    textbox(slide, x+0.15, y+0.42, w-0.3, h-0.55, body, size=9, color=GRAY_500)

def slide_number(slide, n, total=17, dark=False):
    c = RGBColor(0xcc, 0xcc, 0xcc) if not dark else RGBColor(0x44, 0x55, 0x66)
    textbox(slide, 12.5, 7.1, 0.8, 0.3, f"{n:02d} / {total}",
            size=8, color=c, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=SKY_950)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)  # top accent bar

textbox(sl, 0.7, 0.55, 10, 0.35,
        "GRADUATE PROJECT  ·  INDIVIDUAL SUBMISSION  ·  MAY 2026",
        size=9, bold=True, color=SKY_400)

textbox(sl, 0.7, 1.05, 10, 1.0, "FinVerse",
        size=54, bold=True, color=WHITE)
textbox(sl, 0.7, 2.0, 9, 0.9, "Intelligent Personal Finance Platform",
        size=32, bold=True, color=SKY_400)
textbox(sl, 0.7, 3.0, 9, 0.8,
        "A full-stack, multi-module financial analytics application demonstrating real-time\n"
        "account aggregation, AI-powered insights, pattern detection, and predictive forecasting.",
        size=13, color=RGBColor(0xaa, 0xcc, 0xdd))

# tech badges row
badges = ["Next.js 15", "PostgreSQL + Prisma", "Claude AI (Anthropic)",
          "TypeScript", "BullMQ Worker", "RBAC + MFA"]
bx = 0.7
for b in badges:
    w = len(b) * 0.085 + 0.25
    r = rect(sl, bx, 4.1, w, 0.32, fill_rgb=RGBColor(0x15, 0x4a, 0x66))
    r.line.color.rgb = RGBColor(0x30, 0x7a, 0x99)
    r.line.width = Pt(0.5)
    textbox(sl, bx+0.08, 4.12, w-0.1, 0.28, b, size=9, bold=True, color=SKY_400)
    bx += w + 0.12

textbox(sl, 0.7, 6.9, 8, 0.35, "Demo Date: April 29, 2026   ·   Paresh Pandit",
        size=10, color=RGBColor(0x55, 0x77, 0x88))
textbox(sl, 9.5, 6.9, 3.5, 0.35, "17 slides", size=9,
        color=RGBColor(0x33, 0x55, 0x66), align=PP_ALIGN.RIGHT)
slide_number(sl, 1, dark=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is FinVerse
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 4, 0.28, "OVERVIEW", size=9, bold=True, color=SKY_500)
textbox(sl, 0.6, 0.6, 10, 0.65, "What is FinVerse?", size=30, bold=True, color=SKY_950)
textbox(sl, 0.6, 1.3, 11, 0.5,
        "A production-quality financial analytics platform giving individuals a single secure dashboard "
        "over all their financial accounts — with built-in intelligence for detection, forecasting, and advice.",
        size=12, color=GRAY_500)

cards = [
    ("📊  Aggregate", "Unifies checking, savings, credit cards, investment accounts, and loans. Real-time net worth tracking across 7 linked accounts."),
    ("🔍  Detect",    "Automatically identifies recurring subscriptions (≤10% variance algorithm), spending anomalies, and budget overruns using statistical methods."),
    ("🤖  Advise",    "Generates AI-powered health reports, spending forecasts, debt payoff strategies, and tax estimates via Claude Sonnet 4.6."),
]
for i, (t, b) in enumerate(cards):
    section_card(sl, 0.6 + i*4.2, 2.05, 4.0, 2.4, t, b,
                 title_color=SKY_950, bg=SKY_50, border=SKY_100)

textbox(sl, 0.6, 4.65, 3, 0.3, "Core Problem Statement", size=10, bold=True, color=SKY_950)
textbox(sl, 0.6, 4.98, 12, 0.9,
        "Most individuals manage finances across 5–8 disconnected accounts with no holistic view. "
        "FinVerse replaces spreadsheets and app-switching with a unified analytical platform that surfaces "
        "insights automatically — reducing the cognitive load of financial management.",
        size=11, color=GRAY_700)
slide_number(sl, 2)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — System Architecture
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=SKY_50)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 4, 0.28, "TECHNICAL FOUNDATION", size=9, bold=True, color=SKY_500)
textbox(sl, 0.6, 0.6, 10, 0.65, "System Architecture", size=30, bold=True, color=SKY_950)

layers = [
    (RGBColor(0xe0,0xf2,0xfe), SKY_700,  "Presentation Layer", "Next.js 15 App Router · React 18 · Tailwind CSS · Recharts"),
    (EMLD_100,                  EMERALD,  "API Layer",          "Next.js Route Handlers · Zod validation · RBAC middleware · NextAuth v5 JWT"),
    (AMBER_100,                 AMBER,    "Worker / AI Layer",  "BullMQ job queue · Redis · Claude Sonnet 4.6 · Tool-use pipeline"),
    (VIOLT_100,                 VIOLET,   "Data Layer",         "PostgreSQL · Prisma ORM · 27 models · BigInt monetary values"),
    (ROSE_100,                  ROSE_500, "Security Layer",     "AES-256-GCM encryption · MFA (TOTP) · CSP headers · 20+ permissions"),
]
for i, (bg, dot, title, detail) in enumerate(layers):
    y = 1.5 + i * 0.92
    r = rect(sl, 0.6, y, 7.8, 0.78, fill_rgb=bg)
    r.line.color.rgb = dot; r.line.width = Pt(0.5)
    rc = rect(sl, 0.75, y + 0.24, 0.18, 0.18, fill_rgb=dot)
    textbox(sl, 1.1, y + 0.1, 2.5, 0.35, title, size=11, bold=True, color=SKY_950)
    textbox(sl, 1.1, y + 0.42, 7.1, 0.3, detail, size=9, color=GRAY_700)

# Right side stats
stats = [("27", "Prisma Models"), ("20+", "API Routes"), ("4", "RBAC Roles"), ("6", "Packages")]
for i, (v, l) in enumerate(stats):
    xi = 9.0 + (i % 2) * 2.1
    yi = 1.5 + (i // 2) * 1.3
    r = rect(sl, xi, yi, 1.9, 1.1, fill_rgb=WHITE)
    r.line.color.rgb = SKY_100; r.line.width = Pt(0.5)
    textbox(sl, xi, yi+0.1, 1.9, 0.55, v, size=28, bold=True, color=SKY_950, align=PP_ALIGN.CENTER)
    textbox(sl, xi, yi+0.65, 1.9, 0.3, l, size=9, color=GRAY_500, align=PP_ALIGN.CENTER)

textbox(sl, 9.0, 4.2, 4.2, 0.3, "Monorepo Structure", size=10, bold=True, color=SKY_950)
mono_lines = ["apps/web  — Next.js dashboard", "apps/worker  — BullMQ AI worker",
              "packages/db  — Prisma schema + migrations",
              "packages/shared  — RBAC, money, validators"]
for i, l in enumerate(mono_lines):
    textbox(sl, 9.0, 4.55 + i*0.38, 4.2, 0.36, l, size=9, color=GRAY_700)
slide_number(sl, 3)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — April 29 Demo: Core Features
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 8, 0.28, "APRIL 29 DEMO  ·  PART 1 OF 2", size=9, bold=True, color=SKY_500)
textbox(sl, 0.6, 0.6, 10, 0.65, "Core Platform Features", size=30, bold=True, color=SKY_950)
textbox(sl, 0.6, 1.25, 10, 0.35,
        "Seven fully functional modules demonstrated end-to-end with live seeded data.",
        size=12, color=GRAY_500)

features = [
    ("🏠", "Unified Dashboard",  "7 linked accounts — checking ($6.2k), savings ($18.4k), 2 credit cards, Roth IRA, brokerage ($12.5k), auto loan. Net worth at $61.3k."),
    ("📈", "Portfolio Tracker",  "11 holdings across IRA + brokerage. Shows unrealized P&L per position ($13.3k total), asset allocation pie chart, day-change tracking."),
    ("💳", "Credit Score Trends","6-month score trajectory (718→742), utilization breakdown, 3 credit accounts with payment status and next due dates."),
    ("🎯", "Goal Management",    "5 concurrent goals — Emergency Fund (51%), House Down Payment (22% – at risk), Roth IRA Max, Japan Trip (completed ✓), Auto Loan."),
    ("📊", "Budget Management",  "8-category monthly budget. Shopping at 192% overage (MacBook Air – intentional demo spike). Bills at 91%, approaching limit."),
    ("📝", "Net Worth History",  "12-month snapshot trail ($32.6k → $61.3k). Line chart with asset/liability breakdown. Clear upward trend for demo impact."),
]
cols = 3
for i, (icon, title, body) in enumerate(features):
    c = i % cols; r = i // cols
    x = 0.6 + c * 4.2; y = 1.78 + r * 2.3
    bg = rect(sl, x, y, 4.0, 2.15, fill_rgb=SKY_50)
    bg.line.color.rgb = SKY_100; bg.line.width = Pt(0.5)
    textbox(sl, x+0.15, y+0.12, 3.7, 0.3, icon + "  " + title, size=11, bold=True, color=SKY_950)
    textbox(sl, x+0.15, y+0.47, 3.7, 1.55, body, size=9, color=GRAY_700)
slide_number(sl, 4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — April 29 Demo: Advanced Analytics
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=SKY_50)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 8, 0.28, "APRIL 29 DEMO  ·  PART 2 OF 2", size=9, bold=True, color=SKY_500)
textbox(sl, 0.6, 0.6, 10, 0.65, "Advanced Analytics Modules", size=30, bold=True, color=SKY_950)

analytics = [
    ("🤖", "AI Financial Health Report",
     "Claude Sonnet 4.6 generates a structured 0–100 health score with colour-coded Concerns, Strengths, and Recommendations in a 3-column layout. Demo score: 72/100."),
    ("📅", "Spending Heatmap",
     "Month-view calendar where cell darkness maps to daily spend. Rent days and the MacBook purchase are visually distinct peaks. Click any cell for merchant breakdown."),
    ("🔔", "Alert System",
     "4 rule types: budget breach, large transaction, spending threshold, bill due. 2 unread alerts pre-seeded from the MacBook purchase triggering the $500 large-transaction rule."),
    ("🔄", "Recurring Detector",
     "Scans 90 days of transactions, groups by merchant, flags ≥2 charges with ≤10% amount variance. Detected 8 subscriptions totalling $120.94/month."),
    ("📉", "Debt Payoff Calculator",
     "3 debts ($38k total). Real-time amortization slider — drag monthly payment, see payoff timeline and interest savings update instantly. Avalanche vs Snowball comparison."),
    ("🧾", "Tax Estimator",
     "Pre-fills income from transactions and $13.3k capital gains from portfolio. Live 2024 federal bracket + state flat rate. All math client-side, updates on every keystroke."),
]
cols = 2
for i, (icon, title, body) in enumerate(analytics):
    c = i % cols; r = i // cols
    x = 0.6 + c * 6.3; y = 1.5 + r * 1.85
    bg = rect(sl, x, y, 6.0, 1.7, fill_rgb=WHITE)
    bg.line.color.rgb = GRAY_200; bg.line.width = Pt(0.5)
    textbox(sl, x+0.18, y+0.12, 5.6, 0.3, icon + "  " + title, size=11, bold=True, color=SKY_950)
    textbox(sl, x+0.18, y+0.47, 5.6, 1.1, body, size=9, color=GRAY_700)
slide_number(sl, 5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Demo Workflow
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=SKY_950)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 6, 0.28, "USER JOURNEY", size=9, bold=True, color=SKY_400)
textbox(sl, 0.6, 0.6, 11, 0.7, "Demo Workflow — April 29", size=30, bold=True, color=WHITE)

steps = [
    ("1", "Authenticate",         "JWT login (demo@finverse.app). Session via HttpOnly cookie. Optional TOTP MFA for sensitive operations."),
    ("2", "Dashboard Overview",   "Single-page summary: net worth ($61.3k), 7 account balances, recent transactions, portfolio snapshot, and financial health card."),
    ("3", "Portfolio Analysis",   "Navigate to Portfolio → holdings table with per-position P&L, unrealized gains, asset allocation pie chart."),
    ("4", "Generate Health Report","Click 'Generate Report' → BullMQ job → Claude AI analyses 12 financial indicators → structured report in ~8 seconds."),
    ("5", "Spending Intelligence","Heatmap reveals MacBook spike on Day 9. Recurring detector surfaces 8 subscriptions. Budget shows 192% overage in Shopping."),
    ("6", "Debt, Forecast & Tax", "Debt amortization slider in real-time. 30-day forecast at 78% confidence. Tax page pre-fills income + unrealized gains."),
    ("7", "Alert Review",         "2 unread alerts highlighted in amber — MacBook triggered large-transaction rule, Shopping category breached 100% of budget."),
    ("8", "Goal Tracking",        "House Down Payment goal flagged 'at risk' — at current $1,000/mo rate will be $31k short by target date."),
]
for i, (num, title, body) in enumerate(steps):
    c = i % 2; r = i // 2
    x = 0.6 + c * 6.3; y = 1.55 + r * 1.38
    nr = rect(sl, x, y+0.08, 0.35, 0.35, fill_rgb=SKY_500)
    textbox(sl, x+0.01, y+0.1, 0.33, 0.3, num, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, x+0.5, y+0.07, 5.5, 0.28, title, size=11, bold=True, color=WHITE)
    textbox(sl, x+0.5, y+0.38, 5.65, 0.85, body, size=9, color=RGBColor(0xaa,0xcc,0xdd))
slide_number(sl, 6, dark=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — What Was Built: Backend & Data
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 8, 0.28, "WHAT WAS BUILT  ·  PART 1 OF 3", size=9, bold=True, color=SKY_500)
textbox(sl, 0.6, 0.6, 11, 0.7, "Backend & Data Architecture", size=30, bold=True, color=SKY_950)

domains = [
    ("Identity Domain",  "User, Organization, Membership,\nSession, MfaSecret, AuditLog",  "6 models"),
    ("Financial Domain", "FinancialAccount, Transaction,\nHolding, Security, Liability",    "6 models"),
    ("Planning Domain",  "Goal, Budget, BudgetCategory,\nDebtAccount, PayoffStrategy",      "6 models"),
    ("Intelligence",     "Insight, SpendingForecast,\nAlertRule, AlertHistory, CreditScore", "7 models"),
]
for i, (t, d, n) in enumerate(domains):
    x = 0.6 + (i%2)*3.5; y = 1.55 + (i//2)*1.7
    bg = rect(sl, x, y, 3.3, 1.55, fill_rgb=SKY_50)
    bg.line.color.rgb = SKY_100; bg.line.width = Pt(0.5)
    textbox(sl, x+0.15, y+0.1, 3.0, 0.3, t, size=10, bold=True, color=SKY_950)
    textbox(sl, x+0.15, y+0.42, 3.0, 0.75, d, size=8.5, color=GRAY_700)
    textbox(sl, x+0.15, y+1.22, 2.0, 0.25, n, size=8, bold=True, color=SKY_500)

decisions = [
    ("$", "BigInt for All Money",      "Every monetary value stored as integer cents. Eliminates IEEE 754 rounding errors. Explicit centsToDollars() conversion at every JSON boundary."),
    ("🔒", "RBAC with 20+ Permissions","4 roles (Owner/Admin/Member/Viewer) × granular permissions checked on every API route. MFA required for sensitive mutations."),
    ("⚡", "AI via Tool-Use Pipeline", "BullMQ job enqueues analysis. Worker calls Claude with 12 financial tools — get_spending_by_category, get_portfolio_summary, etc. Structured JSON returned."),
    ("🌱", "Production-Ready Seed",    "99 realistic transactions, 11 holdings, 5 goals, 3 debts, 4 alert rules, 3 AI insights — all idempotent and designed to make every feature demo-ready."),
]
for i, (icon, t, b) in enumerate(decisions):
    x = 7.4; y = 1.55 + i * 1.38
    bg = rect(sl, x, y, 5.7, 1.23, fill_rgb=SKY_50)
    bg.line.color.rgb = SKY_100; bg.line.width = Pt(0.5)
    textbox(sl, x+0.15, y+0.1, 5.4, 0.3, icon + "  " + t, size=10, bold=True, color=SKY_950)
    textbox(sl, x+0.15, y+0.42, 5.4, 0.72, b, size=9, color=GRAY_700)
slide_number(sl, 7)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — What Was Built: Algorithms
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=SKY_50)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 8, 0.28, "WHAT WAS BUILT  ·  PART 2 OF 3", size=9, bold=True, color=SKY_500)
textbox(sl, 0.6, 0.6, 10, 0.7, "Algorithms & Detection Logic", size=30, bold=True, color=SKY_950)

algos = [
    ("🔄  Recurring Detector",
     "Groups 90-day transactions by merchant name. For each group with ≥2 charges:\n"
     "  maxDeviation = max(|aᵢ − avg| / avg)\n"
     "  if maxDeviation ≤ 0.10  →  RECURRING\n"
     "  monthlyCost = (avg × 30) / (90 / count)\n"
     "Detected 8 subscriptions in demo. Shared between web API and Claude tool."),
    ("📉  Debt Amortization",
     "Client-side real-time simulation (no API round-trip):\n"
     "  while balance > 0:  interest = balance × (APR/12)\n"
     "                      balance += interest − payment\n"
     "Slider recalculation in <1ms. Compares minimum vs accelerated payoff."),
    ("📊  Spending Forecast",
     "Aggregates 90 days into daily buckets. Computes:\n"
     "  trend = (recent30Avg − older30Avg) / older30Avg\n"
     "  projection[i] = recent30Avg × (1 + trend × i/30)\n"
     "  confidence = 1 − (σ/mean)  [clamped 0.30–0.95]\n"
     "Stored in SpendingForecast table; 78% confidence in demo."),
    ("🧾  Tax Estimation",
     "2024 federal progressive brackets + 22 state flat rates.\n"
     "  taxableIncome = gross − standardDeduction\n"
     "  ltcgTax = bracketTax(income+gains) − bracketTax(income)\n"
     "Pre-fills income from transactions & gains from portfolio.\n"
     "All computation runs on the client — zero server latency."),
]
for i, (title, body) in enumerate(algos):
    c = i % 2; r = i // 2
    x = 0.6 + c * 6.3; y = 1.55 + r * 2.65
    bg = rect(sl, x, y, 6.0, 2.5, fill_rgb=WHITE)
    bg.line.color.rgb = GRAY_200; bg.line.width = Pt(0.5)
    textbox(sl, x+0.18, y+0.12, 5.6, 0.3, title, size=11, bold=True, color=SKY_950)
    textbox(sl, x+0.18, y+0.48, 5.6, 1.85, body, size=8.5, color=GRAY_700)
slide_number(sl, 8)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — What Was Built: UI
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 8, 0.28, "WHAT WAS BUILT  ·  PART 3 OF 3", size=9, bold=True, color=SKY_500)
textbox(sl, 0.6, 0.6, 10, 0.7, "UI Design System & Module Coverage", size=30, bold=True, color=SKY_950)

# Colour swatches
colours = [(SKY_950,"sky-950 Primary"),(SKY_500,"sky-500 Accent"),
           (SKY_100,"sky-100 Surface"),(ROSE_500,"rose Danger"),
           (EMERALD,"emerald Success"),(AMBER,"amber Warning")]
for i,(c,l) in enumerate(colours):
    rx = 0.6 + i*2.1; ry = 1.5
    r = rect(sl, rx, ry, 2.0, 0.42, fill_rgb=c)
    textbox(sl, rx, ry+0.44, 2.0, 0.25, l, size=8, color=GRAY_700, align=PP_ALIGN.CENTER)

textbox(sl, 0.6, 2.35, 4, 0.3, "14 Complete Modules", size=11, bold=True, color=SKY_950)
modules = ["Dashboard","Accounts","Portfolio","Credit",
           "Net Worth","Goals","Budgets","Heatmap",
           "Recurring","Debt","Alerts","Forecast","Tax","Insights"]
for i, m in enumerate(modules):
    c = i % 7; r = i // 7
    x = 0.6 + c * 1.82; y = 2.72 + r * 0.48
    bg = rect(sl, x, y, 1.75, 0.38, fill_rgb=SKY_50)
    bg.line.color.rgb = SKY_100; bg.line.width = Pt(0.5)
    textbox(sl, x+0.1, y+0.05, 1.55, 0.28, m, size=9, bold=False, color=SKY_950)

textbox(sl, 0.6, 3.85, 6, 0.3, "UI Patterns Applied", size=11, bold=True, color=SKY_950)
patterns = [
    "✦  Server components for initial data load (budgets, portfolio)",
    "✦  Client components for interactivity (debt slider, heatmap calendar)",
    "✦  Inline dialogs — no external component file dependencies",
    "✦  Progressive disclosure — expand/collapse on debt and alert cards",
    "✦  Real-time computation — tax and amortization, zero server latency",
]
for i, p in enumerate(patterns):
    textbox(sl, 0.6, 4.18 + i*0.44, 8, 0.38, p, size=10, color=GRAY_700)
slide_number(sl, 9)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Lessons: Technical
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=SKY_50)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 8, 0.28, "LESSONS LEARNED  ·  PART 1 OF 2", size=9, bold=True, color=SKY_500)
textbox(sl, 0.6, 0.6, 11, 0.7, "Technical Lessons", size=30, bold=True, color=SKY_950)

tech_lessons = [
    ("1", "Prisma Schema ≠ Generated Client",
     "Adding @map or new fields requires: run migration, regenerate client, restart server — all three. Skipping any step produces silent P2022 column-not-found errors at runtime, not at compile time."),
    ("2", "pnpm Workspaces Have Non-Obvious Resolution Rules",
     "Each package resolves from its own node_modules first. The Prisma client generated at workspace root was invisible to sub-packages using a local stub — causing TypeScript and runtime divergence."),
    ("3", "JSON Cannot Serialize BigInt or Prisma Decimal",
     "JSON.stringify() throws on BigInt and silently coerces Decimal objects to strings. Every API boundary required explicit Number() and centsToDollars() conversions — a systemic issue across 6 routes."),
    ("4", "Browser Session State Outlives Database Resets",
     "After db:reset, existing JWT cookies still validate (same AUTH_SECRET) but reference stale UUIDs. Dashboard found no org membership → silent redirect loop. Fix: signOut() before re-auth."),
    ("5", "Zod Validator Chaining Order Matters",
     "z.number().default(90).min(30) throws at module load — .default() returns ZodDefault, not ZodNumber. Correct: z.number().min(30).max(365).default(90). Crashed the entire forecast route on startup."),
    ("6", "useEffect Dependencies Cause Race Conditions",
     "Hooks fetching data with [] while using async props (orgId) sometimes executed before the prop was set. Fix: split into two effects — one for orgId resolution, one for data fetching."),
]
for i, (num, title, body) in enumerate(tech_lessons):
    c = i % 2; r = i // 2
    x = 0.6 + c * 6.3; y = 1.55 + r * 1.82
    bg = rect(sl, x, y, 6.0, 1.68, fill_rgb=WHITE)
    bg.line.color.rgb = GRAY_200; bg.line.width = Pt(0.5)
    nr = rect(sl, x+0.15, y+0.17, 0.32, 0.32, fill_rgb=SKY_500)
    textbox(sl, x+0.16, y+0.19, 0.30, 0.28, num, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, x+0.58, y+0.1, 5.25, 0.3, title, size=10, bold=True, color=SKY_950)
    textbox(sl, x+0.58, y+0.43, 5.25, 1.1, body, size=8.5, color=GRAY_700)
slide_number(sl, 10)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Lessons: Design & Product
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 8, 0.28, "LESSONS LEARNED  ·  PART 2 OF 2", size=9, bold=True, color=SKY_500)
textbox(sl, 0.6, 0.6, 11, 0.7, "Design, Product & Architecture Lessons", size=30, bold=True, color=SKY_950)

prod_lessons = [
    ("A", EMERALD, "Seed Data Is a First-Class Feature",
     "Analytical features require dense, realistic data to be meaningful. The final seed — 99 transactions with deliberate peaks, 8 recurring subscriptions, 192% budget spike — tells a story. Sparse seed data is a demo killer."),
    ("B", EMERALD, "Progressive Disclosure Beats Feature Dumping",
     "Early debt UI crammed all details into one scrollable card. Redesigning to summary KPIs + per-card expand/collapse + inline slider kept everything in one mental context and removed the need for modal navigation."),
    ("C", EMERALD, "Design Systems Pay Compound Interest",
     "The first 4 modules used generic shadcn defaults. After a unified sky-950 system, each subsequent module took ~40% less time to style because every decision — color, radius, shadow, spacing — was already resolved."),
    ("D", EMERALD, "Client-Side Computation Beats API Latency for Interactive Tools",
     "Debt amortization and tax estimator originally made API calls on input change. Even 50ms latency made sliders feel laggy. Moving math to pure JS (<1ms) made both features feel instantaneous."),
    ("E", EMERALD, "Error Surfaces Should Diagnose, Not Just Report",
     "The forecast route returned generic 500s for weeks. Adding try/catch with console.error and propagating the actual error message to the client immediately revealed the Zod chaining bug. Opaque 500s are debugging debt."),
    ("F", EMERALD, "Architecture Decisions Have Long Shadows",
     "BigInt cents prevents precision bugs but requires explicit conversion everywhere. pnpm monorepo enables code sharing but adds Prisma resolution complexity. Each tradeoff compounds — early decisions are load-bearing walls."),
]
for i, (letter, color, title, body) in enumerate(prod_lessons):
    c = i % 2; r = i // 2
    x = 0.6 + c * 6.3; y = 1.55 + r * 1.82
    bg = rect(sl, x, y, 6.0, 1.68, fill_rgb=SKY_50)
    bg.line.color.rgb = SKY_100; bg.line.width = Pt(0.5)
    nr = rect(sl, x+0.15, y+0.17, 0.32, 0.32, fill_rgb=color)
    textbox(sl, x+0.16, y+0.19, 0.30, 0.28, letter, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, x+0.58, y+0.1, 5.25, 0.3, title, size=10, bold=True, color=SKY_950)
    textbox(sl, x+0.58, y+0.43, 5.25, 1.1, body, size=8.5, color=GRAY_700)
slide_number(sl, 11)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Big Data: Current Scale & Bottlenecks
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=SKY_950)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 8, 0.28, "BIG DATA ANALYSIS  ·  PART 1 OF 3", size=9, bold=True, color=SKY_400)
textbox(sl, 0.6, 0.6, 11, 0.7, "Current Scale & Breaking Points", size=30, bold=True, color=WHITE)

# Current stats
cur_stats = [("99", "transactions\nper user"), ("11", "portfolio\nholdings"),
             ("12", "net worth\nsnapshots"), ("~2ms", "avg query\ntime")]
for i, (v, l) in enumerate(cur_stats):
    x = 0.6 + i*3.1; y = 1.5
    r = rect(sl, x, y, 2.9, 1.3, fill_rgb=RGBColor(0x0c,0x3a,0x55))
    r.line.color.rgb = RGBColor(0x1a,0x5a,0x80); r.line.width = Pt(0.5)
    textbox(sl, x, y+0.1, 2.9, 0.6, v, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, x, y+0.72, 2.9, 0.5, l, size=9, color=RGBColor(0x88,0xaa,0xbb), align=PP_ALIGN.CENTER)

textbox(sl, 0.6, 3.05, 12, 0.35,
        "At 10M users × 500 transactions/year → 5 billion rows. Four critical bottlenecks:",
        size=11, color=RGBColor(0xaa,0xcc,0xdd))

bottlenecks = [
    (ROSE_500,  "🔴  Full-Table Scans for Recurring Detection",
     "90-day merchant GROUP BY requires scanning every row per org. No index helps a GROUP BY merchantName across 5B rows. Estimated query time: >30 seconds."),
    (AMBER,     "🟡  Budget Aggregation on Every Page Load",
     "Budget spent amounts are recomputed from raw transactions on each dashboard load. At scale, this cannot sustain real-time ingest without pre-aggregated materialized views."),
    (SKY_400,   "🔵  Per-User AI Insight Generation",
     "Each health report uses ~3,000 input tokens. At 10M users × weekly reports = 1.56T tokens/year ≈ $2M+ at current Claude API pricing. Requires batching, caching, and model distillation."),
    (VIOLET,    "🟣  Synchronous Forecast Computation",
     "Linear regression computed on-demand per request. At scale, forecasts must be pre-computed in batch nightly and cached — not recomputed per page load."),
]
for i, (color, title, body) in enumerate(bottlenecks):
    x = 0.6 + (i%2)*6.3; y = 3.55 + (i//2)*1.65
    r = rect(sl, x, y, 6.0, 1.52, fill_rgb=RGBColor(0x0c,0x2a,0x3f))
    r.line.color.rgb = color; r.line.width = Pt(1.2)
    textbox(sl, x+0.18, y+0.1, 5.6, 0.3, title, size=10, bold=True, color=WHITE)
    textbox(sl, x+0.18, y+0.43, 5.6, 0.95, body, size=8.5, color=RGBColor(0xaa,0xcc,0xdd))
slide_number(sl, 12, dark=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Big Data: Infrastructure Solutions
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=SKY_50)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 8, 0.28, "BIG DATA ANALYSIS  ·  PART 2 OF 3", size=9, bold=True, color=SKY_500)
textbox(sl, 0.6, 0.6, 11, 0.7, "Infrastructure Solutions at Scale", size=30, bold=True, color=SKY_950)

solutions_left = [
    (SKY_700,  "Apache Kafka — Real-Time Event Streaming",
     "Transaction events published to finverse.transactions topic. Alert evaluation, budget aggregation, and anomaly detection subscribe as independent consumer groups — eliminates polling-based BullMQ approach."),
    (EMERALD,  "Apache Spark — Batch Analytics",
     "Nightly Spark jobs: recurring detection as distributed GroupBy, spending pattern clustering via k-means, forecast pre-computation via distributed linear regression across partitioned user data."),
    (AMBER,    "Data Warehouse (Snowflake / BigQuery)",
     "Separate OLAP tier from PostgreSQL OLTP. Transaction history, net worth snapshots, and portfolio valuations land in columnar storage optimised for aggregation — not live transaction rows."),
    (ROSE_500, "Delta Lake / Parquet — Time-Series Storage",
     "Net worth snapshots stored in Parquet partitioned by (user_id, year, month). ACID transactions via Delta Lake. Reduces net worth history from O(rows) to O(partitions)."),
]
solutions_right = [
    ("Redis Distributed Cache",
     "Pre-computed results with TTL: budget aggregations (5 min), recurring subscriptions (24 hr), forecast data (12 hr). Cache invalidation triggered via Kafka on new transaction ingest."),
    ("PostgreSQL Read Replicas",
     "Streaming replication for reads. Write operations go to primary; dashboard reads, portfolio queries, and credit lookups route to read replicas — eliminating read/write contention."),
    ("Materialized Views for Aggregations",
     "Budget category spending, monthly net worth pre-computed as PostgreSQL materialized views. Refreshed incrementally on transaction insert via Kafka consumer — not on every page load."),
    ("AI Cost Optimization",
     "Embed spending patterns as 768-dim vectors. Batch similar users together. Use Claude only for final narrative generation on pre-computed summaries — reducing per-user token consumption by ~80%."),
]
for i, (color, title, body) in enumerate(solutions_left):
    y = 1.55 + i * 1.38
    r = rect(sl, 0.6, y, 6.0, 1.25, fill_rgb=WHITE)
    r.line.color.rgb = color; r.line.width = Pt(1.5)
    textbox(sl, 0.75, y+0.1, 5.7, 0.3, title, size=10, bold=True, color=SKY_950)
    textbox(sl, 0.75, y+0.43, 5.7, 0.72, body, size=8.5, color=GRAY_700)

for i, (title, body) in enumerate(solutions_right):
    y = 1.55 + i * 1.38
    bg = rect(sl, 6.9, y, 6.0, 1.25, fill_rgb=SKY_50)
    bg.line.color.rgb = SKY_100; bg.line.width = Pt(0.5)
    textbox(sl, 7.05, y+0.1, 5.7, 0.3, title, size=10, bold=True, color=SKY_950)
    textbox(sl, 7.05, y+0.43, 5.7, 0.72, body, size=8.5, color=GRAY_700)
slide_number(sl, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Big Data: Transformative Opportunities
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 8, 0.28, "BIG DATA ANALYSIS  ·  PART 3 OF 3", size=9, bold=True, color=SKY_500)
textbox(sl, 0.6, 0.6, 11, 0.7, "Transformative Opportunities at Population Scale", size=30, bold=True, color=SKY_950)
textbox(sl, 0.6, 1.32, 11, 0.35,
        "Big data doesn't just solve the scaling problem — it enables entirely new product capabilities impossible with single-user data.",
        size=11, color=GRAY_500)

opportunities = [
    ("🏘️", "Peer Benchmarking",
     "'Your food spending is 34% above similar households.' Cross-user analytics via k-anonymity / differential privacy transforms individual dashboards into contextualised intelligence.",
     "Spark aggregations + census cohort modelling"),
    ("🔍", "ML Fraud Detection",
     "Isolation Forest over transaction vectors (merchant, amount, time) trained on 5B+ labelled transactions. Detects anomalous charges in real time — impossible without population-scale training data.",
     "Kafka streams + Seldon/BentoML model serving"),
    ("📉", "Predictive Credit Modelling",
     "Feature-engineer 40+ signals from spending patterns and payment timing. Train gradient-boosted classifier to predict credit score change 60 days in advance — actionable, not retrospective.",
     "Feature store + XGBoost pipeline"),
    ("📊", "Market Correlation Engine",
     "Correlate portfolio holdings with macro indicators (CPI, Fed rate, sector ETF flows). Proactive: 'Your NVDA position correlates 0.84 with NASDAQ — consider rebalancing.'",
     "Time-series DB + financial data APIs"),
    ("🤖", "Personalised AI Distillation",
     "Fine-tune a smaller model (Haiku-class) on FinVerse financial reasoning using RLHF from user feedback signals. Reduces inference cost 10× vs Claude Sonnet while maintaining domain accuracy.",
     "SFT pipeline + feedback telemetry"),
    ("🔔", "Proactive Alert Intelligence",
     "'Your gym membership has been unused for 47 days — 89% of users in this situation cancel within 60 days.' Rule-based thresholds replaced with ML-driven cohort behaviour analysis.",
     "Behavioural cohort analysis + Kafka streaming"),
]
for i, (icon, title, body, tech) in enumerate(opportunities):
    c = i % 3; r = i // 3
    x = 0.6 + c * 4.2; y = 1.88 + r * 2.55
    bg = rect(sl, x, y, 4.0, 2.4, fill_rgb=SKY_50)
    bg.line.color.rgb = SKY_100; bg.line.width = Pt(0.5)
    textbox(sl, x+0.15, y+0.12, 3.7, 0.28, icon + "  " + title, size=10, bold=True, color=SKY_950)
    textbox(sl, x+0.15, y+0.44, 3.7, 1.38, body, size=8.5, color=GRAY_700)
    textbox(sl, x+0.15, y+2.08, 3.7, 0.25, "→ " + tech, size=7.5, bold=True, color=SKY_500)
slide_number(sl, 14)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill_rgb=SKY_950)
rect(sl, 0, 0, 13.33, 0.06, fill_rgb=SKY_500)

textbox(sl, 0.6, 0.25, 4, 0.28, "CONCLUSION", size=9, bold=True, color=SKY_400)
textbox(sl, 0.6, 0.6, 8, 0.7, "FinVerse: From Prototype to\nProduction-Grade Platform", size=28, bold=True, color=WHITE)

highlights = [
    ("🏗️", "Full-Stack Engineering Depth",
     "Built end-to-end: database schema, API architecture, authentication, background processing, AI integration, and a 14-module UI — demonstrating production-quality engineering across the entire stack."),
    ("🧠", "Algorithms Applied to Real Data",
     "Recurring detection, linear regression forecasting, amortization simulation, and tax bracket calculation — all implemented from first principles, grounded in real financial mathematics."),
    ("🔬", "Big Data Readiness Analysis",
     "Systematically identified where the synchronous single-node architecture fails at population scale, and mapped specific infrastructure solutions (Kafka, Spark, Delta Lake, Redis) to each bottleneck."),
]
for i, (icon, title, body) in enumerate(highlights):
    y = 1.6 + i * 1.52
    r = rect(sl, 0.6, y, 7.8, 1.38, fill_rgb=RGBColor(0x0c,0x3a,0x55))
    r.line.color.rgb = RGBColor(0x1a,0x5a,0x80); r.line.width = Pt(0.5)
    textbox(sl, 0.9, y+0.1, 0.55, 0.4, icon, size=20, color=WHITE)
    textbox(sl, 1.6, y+0.1, 6.6, 0.3, title, size=11, bold=True, color=WHITE)
    textbox(sl, 1.6, y+0.43, 6.6, 0.82, body, size=9, color=RGBColor(0xaa,0xcc,0xdd))

# Key metrics
metrics = [("14","Modules"),("27","DB Models"),("20+","API Routes"),("4","Algorithms")]
for i,(v,l) in enumerate(metrics):
    x = 9.1; y = 1.6 + i*1.3
    r = rect(sl, x, y, 3.9, 1.15, fill_rgb=RGBColor(0x0c,0x2a,0x3f))
    r.line.color.rgb = RGBColor(0x1a,0x5a,0x80); r.line.width = Pt(0.5)
    textbox(sl, x, y+0.05, 3.9, 0.6, v, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, x, y+0.72, 3.9, 0.35, l, size=10, color=SKY_400, align=PP_ALIGN.CENTER)

textbox(sl, 0.6, 6.85, 8, 0.3, "Paresh Pandit  ·  Graduate Submission  ·  May 2026",
        size=9, color=RGBColor(0x44,0x66,0x77))
textbox(sl, 9.0, 6.85, 4.0, 0.3, "demo@finverse.app  /  Demo@2024",
        size=9, color=RGBColor(0x33,0x55,0x66), align=PP_ALIGN.RIGHT)
slide_number(sl, 15, dark=True)

# ── save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "FinVerse_Demo_Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
